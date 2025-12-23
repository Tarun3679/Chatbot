#!/usr/bin/env python3
"""
LibreOffice OOM Test File Generator

Generates PowerPoint and Excel files of varying complexity to stress-test
LibreOffice PDF conversion and reproduce OOM errors.

Usage:
    python test_file_generator.py [--output-dir ./test_files] [--level medium]
    
Levels:
    light   - Quick tests, small files
    medium  - Moderate stress test
    heavy   - High memory usage
    extreme - Maximum stress (may cause OOM on low-memory systems)
"""

import argparse
import io
import random
import sys
from pathlib import Path

# Check for required dependencies
MISSING_DEPS = []

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RgbColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    MISSING_DEPS.append("python-pptx")

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
    from openpyxl.chart import BarChart, LineChart, PieChart, Reference
    from openpyxl.utils import get_column_letter
except ImportError:
    MISSING_DEPS.append("openpyxl")

try:
    from PIL import Image
    import numpy as np
except ImportError:
    MISSING_DEPS.extend(["pillow", "numpy"])


def check_dependencies():
    """Check if all required dependencies are installed."""
    if MISSING_DEPS:
        print("=" * 60)
        print("MISSING DEPENDENCIES")
        print("=" * 60)
        print(f"Please install the following packages:\n")
        print(f"    pip install {' '.join(MISSING_DEPS)}")
        print("\nOr install all at once:")
        print("    pip install python-pptx openpyxl pillow numpy")
        print("=" * 60)
        sys.exit(1)


def generate_random_text(word_count: int = 100) -> str:
    """Generate random lorem ipsum-style text."""
    words = [
        'lorem', 'ipsum', 'dolor', 'sit', 'amet', 'consectetur',
        'adipiscing', 'elit', 'sed', 'do', 'eiusmod', 'tempor',
        'incididunt', 'ut', 'labore', 'et', 'dolore', 'magna', 'aliqua',
        'enim', 'ad', 'minim', 'veniam', 'quis', 'nostrud', 'exercitation',
        'ullamco', 'laboris', 'nisi', 'aliquip', 'ex', 'ea', 'commodo',
        'consequat', 'duis', 'aute', 'irure', 'in', 'reprehenderit',
        'voluptate', 'velit', 'esse', 'cillum', 'fugiat', 'nulla', 'pariatur',
        'excepteur', 'sint', 'occaecat', 'cupidatat', 'non', 'proident',
        'sunt', 'culpa', 'qui', 'officia', 'deserunt', 'mollit', 'anim'
    ]
    return ' '.join(random.choices(words, k=word_count)).capitalize()


def generate_random_image(width: int = 800, height: int = 600, complexity: str = "medium") -> bytes:
    """
    Generate a random image with varying complexity.
    
    Args:
        width: Image width in pixels
        height: Image height in pixels
        complexity: 'low' (solid color), 'medium' (gradients), 'high' (noise)
    
    Returns:
        PNG image as bytes
    """
    if complexity == "low":
        # Solid color with some shapes
        img_array = np.full((height, width, 3), 
                           fill_value=random.randint(100, 255), 
                           dtype=np.uint8)
    elif complexity == "high":
        # Full random noise (highest entropy, larger file)
        img_array = np.random.randint(0, 255, (height, width, 3), dtype=np.uint8)
    else:  # medium
        # Gradient with noise
        img_array = np.zeros((height, width, 3), dtype=np.uint8)
        for i in range(3):
            gradient = np.linspace(0, 255, width, dtype=np.uint8)
            img_array[:, :, i] = np.tile(gradient, (height, 1))
            img_array[:, :, i] += np.random.randint(0, 50, (height, width), dtype=np.uint8)
    
    img = Image.fromarray(img_array)
    buffer = io.BytesIO()
    img.save(buffer, format='PNG', optimize=False)
    return buffer.getvalue()


def create_heavy_powerpoint(
    num_slides: int = 50,
    images_per_slide: int = 2,
    text_boxes_per_slide: int = 5,
    image_size: tuple = (1024, 768),
    include_tables: bool = True,
    include_charts: bool = True,
    verbose: bool = True,
) -> bytes:
    """
    Create a memory-intensive PowerPoint file.
    
    Args:
        num_slides: Number of slides to create
        images_per_slide: Number of images per slide (major memory driver)
        text_boxes_per_slide: Number of text boxes per slide
        image_size: (width, height) of generated images
        include_tables: Whether to include tables
        include_charts: Whether to include charts
        verbose: Print progress
    
    Returns:
        PowerPoint file as bytes
    """
    if verbose:
        print(f"    Creating PowerPoint: {num_slides} slides, {images_per_slide} images/slide...")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]  # Blank slide
    
    for slide_num in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.7))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = f"Slide {slide_num + 1}: {generate_random_text(8)}"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        
        # Text boxes
        for i in range(text_boxes_per_slide):
            left = Inches(0.5 + (i % 2) * 6.2)
            top = Inches(1.0 + (i // 2) * 1.3)
            
            text_box = slide.shapes.add_textbox(left, top, Inches(5.8), Inches(1.1))
            tf = text_box.text_frame
            tf.word_wrap = True
            
            para = tf.paragraphs[0]
            para.text = generate_random_text(40)
            para.font.size = Pt(11)
        
        # Images (the main memory driver)
        for i in range(images_per_slide):
            img_bytes = generate_random_image(
                width=image_size[0], 
                height=image_size[1],
                complexity="high" if i % 2 == 0 else "medium"
            )
            img_stream = io.BytesIO(img_bytes)
            
            left = Inches(0.3 + (i % 2) * 6.5)
            top = Inches(4.0 + (i // 2) * 1.8)
            
            try:
                slide.shapes.add_picture(img_stream, left, top, width=Inches(6))
            except Exception:
                pass  # Skip if image fails
        
        # Tables
        if include_tables and slide_num % 4 == 0:
            rows, cols = 6, 5
            try:
                table_shape = slide.shapes.add_table(
                    rows, cols, 
                    Inches(1), Inches(2.8), 
                    Inches(10), Inches(1.8)
                )
                table = table_shape.table
                
                for row_idx in range(rows):
                    for col_idx in range(cols):
                        cell = table.cell(row_idx, col_idx)
                        if row_idx == 0:
                            cell.text = f"Header {col_idx + 1}"
                        else:
                            cell.text = f"{random.randint(100, 9999)}"
            except Exception:
                pass
        
        # Progress indicator
        if verbose and (slide_num + 1) % 10 == 0:
            print(f"      Progress: {slide_num + 1}/{num_slides} slides")
    
    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    if verbose:
        size_mb = len(buffer.getvalue()) / (1024 * 1024)
        print(f"    PowerPoint created: {size_mb:.2f} MB")
    
    return buffer.getvalue()


def create_heavy_excel(
    num_sheets: int = 10,
    rows_per_sheet: int = 10000,
    cols_per_sheet: int = 20,
    include_charts: bool = True,
    include_formatting: bool = True,
    include_formulas: bool = True,
    verbose: bool = True,
) -> bytes:
    """
    Create a memory-intensive Excel file.
    
    Args:
        num_sheets: Number of worksheets
        rows_per_sheet: Rows per sheet
        cols_per_sheet: Columns per sheet
        include_charts: Whether to add charts
        include_formatting: Whether to apply cell formatting
        include_formulas: Whether to add formulas
        verbose: Print progress
    
    Returns:
        Excel file as bytes
    """
    if verbose:
        print(f"    Creating Excel: {num_sheets} sheets, {rows_per_sheet} rows each...")
    
    wb = Workbook()
    default_sheet = wb.active
    
    # Styles
    header_font = Font(bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    alt_fill = PatternFill(start_color="D9E2F3", end_color="D9E2F3", fill_type="solid")
    border = Border(
        left=Side(style='thin', color='B4B4B4'),
        right=Side(style='thin', color='B4B4B4'),
        top=Side(style='thin', color='B4B4B4'),
        bottom=Side(style='thin', color='B4B4B4')
    )
    
    for sheet_num in range(num_sheets):
        if sheet_num == 0:
            ws = default_sheet
            ws.title = f"DataSheet_{sheet_num + 1}"
        else:
            ws = wb.create_sheet(title=f"DataSheet_{sheet_num + 1}")
        
        # Headers
        for col_idx in range(1, cols_per_sheet + 1):
            cell = ws.cell(row=1, column=col_idx)
            cell.value = f"Column_{get_column_letter(col_idx)}"
            if include_formatting:
                cell.font = header_font
                cell.fill = header_fill
                cell.border = border
                cell.alignment = Alignment(horizontal='center')
        
        # Data rows
        for row_idx in range(2, rows_per_sheet + 2):
            for col_idx in range(1, cols_per_sheet + 1):
                cell = ws.cell(row=row_idx, column=col_idx)
                
                # Vary data types
                col_type = col_idx % 5
                if col_type == 0:
                    cell.value = generate_random_text(3)
                elif col_type == 1:
                    cell.value = random.uniform(-1000, 10000)
                elif col_type == 2:
                    cell.value = random.randint(-10000, 100000)
                elif col_type == 3:
                    cell.value = f"{random.randint(1, 12)}/{random.randint(1, 28)}/2024"
                else:
                    cell.value = random.choice([True, False, None, "N/A"])
                
                # Formatting
                if include_formatting:
                    cell.border = border
                    if row_idx % 2 == 0:
                        cell.fill = alt_fill
            
            # Progress for large sheets
            if verbose and row_idx % 5000 == 0:
                print(f"      Sheet {sheet_num + 1}: {row_idx}/{rows_per_sheet} rows")
        
        # Formulas
        if include_formulas:
            # Add SUM formulas in last column
            for row_idx in range(2, min(rows_per_sheet + 2, 500)):
                ws.cell(row=row_idx, column=cols_per_sheet).value = \
                    f"=SUM(B{row_idx}:{get_column_letter(cols_per_sheet - 1)}{row_idx})"
            
            # Add summary row
            summary_row = rows_per_sheet + 3
            for col_idx in range(2, cols_per_sheet + 1):
                col_letter = get_column_letter(col_idx)
                ws.cell(row=summary_row, column=col_idx).value = \
                    f"=AVERAGE({col_letter}2:{col_letter}{rows_per_sheet + 1})"
        
        # Charts
        if include_charts and sheet_num % 2 == 0:
            try:
                # Bar chart
                chart = BarChart()
                chart.title = f"Data Analysis - Sheet {sheet_num + 1}"
                chart.x_axis.title = "Row"
                chart.y_axis.title = "Value"
                
                data = Reference(ws, min_col=2, min_row=1, max_col=5, max_row=min(52, rows_per_sheet))
                categories = Reference(ws, min_col=1, min_row=2, max_row=min(52, rows_per_sheet))
                
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(categories)
                chart.width = 18
                chart.height = 10
                
                ws.add_chart(chart, f"{get_column_letter(cols_per_sheet + 2)}2")
            except Exception:
                pass  # Skip chart on error
        
        if verbose:
            print(f"      Completed sheet {sheet_num + 1}/{num_sheets}")
    
    # Save to bytes
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    
    if verbose:
        size_mb = len(buffer.getvalue()) / (1024 * 1024)
        print(f"    Excel created: {size_mb:.2f} MB")
    
    return buffer.getvalue()


# Test file configurations by stress level
TEST_CONFIGS = {
    "light": [
        {"name": "pptx_light_01.pptx", "type": "pptx", 
         "params": {"num_slides": 5, "images_per_slide": 0, "text_boxes_per_slide": 2}},
        {"name": "pptx_light_02.pptx", "type": "pptx",
         "params": {"num_slides": 10, "images_per_slide": 1, "text_boxes_per_slide": 2, "image_size": (640, 480)}},
        {"name": "xlsx_light_01.xlsx", "type": "xlsx",
         "params": {"num_sheets": 2, "rows_per_sheet": 500, "cols_per_sheet": 10}},
        {"name": "xlsx_light_02.xlsx", "type": "xlsx",
         "params": {"num_sheets": 3, "rows_per_sheet": 1000, "cols_per_sheet": 15}},
    ],
    "medium": [
        {"name": "pptx_medium_01.pptx", "type": "pptx",
         "params": {"num_slides": 20, "images_per_slide": 1, "text_boxes_per_slide": 3}},
        {"name": "pptx_medium_02.pptx", "type": "pptx",
         "params": {"num_slides": 30, "images_per_slide": 2, "text_boxes_per_slide": 4}},
        {"name": "pptx_medium_03.pptx", "type": "pptx",
         "params": {"num_slides": 25, "images_per_slide": 2, "text_boxes_per_slide": 3, "image_size": (1280, 720)}},
        {"name": "xlsx_medium_01.xlsx", "type": "xlsx",
         "params": {"num_sheets": 5, "rows_per_sheet": 5000, "cols_per_sheet": 20}},
        {"name": "xlsx_medium_02.xlsx", "type": "xlsx",
         "params": {"num_sheets": 8, "rows_per_sheet": 8000, "cols_per_sheet": 25}},
    ],
    "heavy": [
        {"name": "pptx_heavy_01.pptx", "type": "pptx",
         "params": {"num_slides": 50, "images_per_slide": 3, "text_boxes_per_slide": 5}},
        {"name": "pptx_heavy_02.pptx", "type": "pptx",
         "params": {"num_slides": 75, "images_per_slide": 2, "text_boxes_per_slide": 4, "image_size": (1920, 1080)}},
        {"name": "pptx_heavy_03.pptx", "type": "pptx",
         "params": {"num_slides": 60, "images_per_slide": 4, "text_boxes_per_slide": 6}},
        {"name": "xlsx_heavy_01.xlsx", "type": "xlsx",
         "params": {"num_sheets": 10, "rows_per_sheet": 20000, "cols_per_sheet": 30}},
        {"name": "xlsx_heavy_02.xlsx", "type": "xlsx",
         "params": {"num_sheets": 15, "rows_per_sheet": 15000, "cols_per_sheet": 35}},
    ],
    "extreme": [
        {"name": "pptx_extreme_01.pptx", "type": "pptx",
         "params": {"num_slides": 100, "images_per_slide": 4, "text_boxes_per_slide": 6, "image_size": (1920, 1080)}},
        {"name": "pptx_extreme_02.pptx", "type": "pptx",
         "params": {"num_slides": 150, "images_per_slide": 3, "text_boxes_per_slide": 5, "image_size": (2560, 1440)}},
        {"name": "xlsx_extreme_01.xlsx", "type": "xlsx",
         "params": {"num_sheets": 20, "rows_per_sheet": 50000, "cols_per_sheet": 40}},
        {"name": "xlsx_extreme_02.xlsx", "type": "xlsx",
         "params": {"num_sheets": 25, "rows_per_sheet": 40000, "cols_per_sheet": 50}},
    ],
}


def create_test_suite(output_dir: str = "./test_files", level: str = "medium", verbose: bool = True) -> list:
    """
    Create a suite of test files at the specified stress level.
    
    Args:
        output_dir: Directory to save test files
        level: Stress level ('light', 'medium', 'heavy', 'extreme')
        verbose: Print progress
    
    Returns:
        List of created file info dicts
    """
    check_dependencies()
    
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    configs = TEST_CONFIGS.get(level, TEST_CONFIGS["medium"])
    
    if verbose:
        print("=" * 60)
        print(f"CREATING TEST FILES - Level: {level.upper()}")
        print(f"Output directory: {output_path.absolute()}")
        print("=" * 60)
    
    created_files = []
    
    for i, config in enumerate(configs, 1):
        if verbose:
            print(f"\n[{i}/{len(configs)}] Creating {config['name']}...")
        
        try:
            if config["type"] == "pptx":
                file_bytes = create_heavy_powerpoint(**config["params"], verbose=verbose)
            else:
                file_bytes = create_heavy_excel(**config["params"], verbose=verbose)
            
            filepath = output_path / config["name"]
            filepath.write_bytes(file_bytes)
            
            size_mb = len(file_bytes) / (1024 * 1024)
            created_files.append({
                "path": str(filepath),
                "name": config["name"],
                "type": config["type"],
                "size_mb": size_mb,
                "params": config["params"]
            })
            
            if verbose:
                print(f"    ✓ Saved: {filepath.name} ({size_mb:.2f} MB)")
            
        except Exception as e:
            if verbose:
                print(f"    ✗ Failed: {e}")
    
    if verbose:
        print("\n" + "=" * 60)
        print("SUMMARY")
        print("=" * 60)
        total_size = sum(f["size_mb"] for f in created_files)
        print(f"Files created: {len(created_files)}")
        print(f"Total size: {total_size:.2f} MB")
        print(f"Location: {output_path.absolute()}")
        print("=" * 60)
    
    return created_files


def main():
    parser = argparse.ArgumentParser(
        description="Generate test files for LibreOffice OOM stress testing",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Stress Levels:
  light    - Small files for quick tests (5-10 slides, 500-1000 rows)
  medium   - Moderate stress test (20-30 slides, 5000-8000 rows)
  heavy    - High memory usage (50-75 slides, 15000-20000 rows)
  extreme  - Maximum stress (100-150 slides, 40000-50000 rows)
             WARNING: May cause OOM on systems with <8GB RAM

Examples:
  %(prog)s --level light
  %(prog)s --level heavy --output-dir /tmp/test_files
  %(prog)s --level extreme  # Use with caution!
        """
    )
    
    parser.add_argument(
        "--output-dir", "-o",
        default="./test_files",
        help="Output directory for test files (default: ./test_files)"
    )
    parser.add_argument(
        "--level", "-l",
        choices=["light", "medium", "heavy", "extreme"],
        default="medium",
        help="Stress level (default: medium)"
    )
    parser.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="Suppress progress output"
    )
    
    args = parser.parse_args()
    
    create_test_suite(
        output_dir=args.output_dir,
        level=args.level,
        verbose=not args.quiet
    )


if __name__ == "__main__":
    main()
